from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = PROJECT_ROOT / "DataDictionary_TradeFairBook.pptx"

HEADERS = ["SR.No", "Column Name", "Data type", "Size", "Constrain", "Description"]


def row(sr_no, column_name, data_type, size, constrain, description):
    return [sr_no, column_name, data_type, size, constrain, description]


DATA_DICTIONARY = {
    "Users Collection (users)": [
        row("1", "_id", "ObjectId", "24 hex chars", "PK, Auto-generated", "Unique user identifier"),
        row("2", "fullname", "String", "Variable", "Required", "User full name"),
        row("3", "company", "String", "Variable", "Required", "Company or organization name"),
        row("4", "email", "String", "Variable", "Required, Unique", "User login email"),
        row("5", "phone", "String", "Variable", "Required", "User phone number"),
        row("6", "password", "String", "About 60 chars", "Required", "Bcrypt password hash"),
        row("7", "role", "String", "Up to 11 chars", "Enum: USER,ADMIN,SUPER_ADMIN; Default USER", "Access control role"),
        row("8", "createdAt", "Date", "8 bytes", "Auto (timestamps)", "Record creation time"),
        row("9", "updatedAt", "Date", "8 bytes", "Auto (timestamps)", "Record last update time"),
    ],
    "Domes Collection (domes)": [
        row("1", "_id", "ObjectId", "24 hex chars", "PK, Auto-generated", "Unique dome identifier"),
        row("2", "domeName", "String", "Variable", "Required", "Dome name"),
        row("3", "location", "String", "Variable", "Required", "Dome location"),
        row("4", "description", "String", "Variable", "Optional", "Dome description"),
        row("5", "image", "String", "Variable", "Optional", "Dome image URL or path"),
        row("6", "status", "String", "Up to 8 chars", "Enum: ACTIVE,INACTIVE; Default ACTIVE", "Dome status"),
        row("7", "createdAt", "Date", "8 bytes", "Auto (timestamps)", "Record creation time"),
        row("8", "updatedAt", "Date", "8 bytes", "Auto (timestamps)", "Record last update time"),
    ],
    "Stalls Collection (stalls)": [
        row("1", "_id", "ObjectId", "24 hex chars", "PK, Auto-generated", "Unique stall identifier"),
        row("2", "stallNumber", "String", "Variable", "Required, Trim, Unique with dome", "Stall number or code"),
        row("3", "dome", "ObjectId", "24 hex chars", "Required, FK -> domes._id", "Dome reference"),
        row("4", "side", "String", "Up to 6 chars", "Required, Enum: LEFT,RIGHT,TOP,BOTTOM,CENTER", "Stall side in layout"),
        row("5", "price", "Number", "8 bytes", "Required, Min 0", "Stall price"),
        row("6", "status", "String", "Up to 9 chars", "Enum: AVAILABLE,BOOKED,HOLD,BLOCKED; Default AVAILABLE", "Availability status"),
        row("7", "centerSpacing", "String", "Up to 10 chars", "Enum: with-space,no-space; Default with-space", "Center layout spacing flag"),
        row("8", "createdAt", "Date", "8 bytes", "Auto (timestamps)", "Record creation time"),
        row("9", "updatedAt", "Date", "8 bytes", "Auto (timestamps)", "Record last update time"),
    ],
    "Materials Collection (materials)": [
        row("1", "_id", "ObjectId", "24 hex chars", "PK, Auto-generated", "Unique material identifier"),
        row("2", "dome", "ObjectId", "24 hex chars", "Required, FK -> domes._id", "Dome reference"),
        row("3", "name", "String", "Variable", "Required, Trim, Unique with dome", "Material name"),
        row("4", "price", "Number", "8 bytes", "Required, Min 0", "Material price"),
        row("5", "description", "String", "Variable", "Default empty, Trim", "Material description"),
        row("6", "isActive", "Boolean", "1 byte", "Default true", "Material active flag"),
        row("7", "createdAt", "Date", "8 bytes", "Auto (timestamps)", "Record creation time"),
        row("8", "updatedAt", "Date", "8 bytes", "Auto (timestamps)", "Record last update time"),
    ],
    "AadhaarVerification Collection (aadhaarverifications)": [
        row("1", "_id", "ObjectId", "24 hex chars", "PK, Auto-generated", "Unique verification identifier"),
        row("2", "user", "ObjectId", "24 hex chars", "Required, FK -> users._id", "User reference"),
        row("3", "aadhaarName", "String", "Variable", "Required, Trim", "Name as on Aadhaar"),
        row("4", "aadhaarNumber", "String", "Encrypted variable text", "Required, select:false, API expects 12 digits", "Encrypted Aadhaar number"),
        row("5", "aadhaarImage", "String", "Variable", "Required, Trim", "Uploaded image path"),
        row("6", "verified", "Boolean", "1 byte", "Default false", "Verification status"),
        row("7", "submittedAt", "Date", "8 bytes", "Default Date.now", "Submission timestamp"),
        row("8", "createdAt", "Date", "8 bytes", "Auto (timestamps)", "Record creation time"),
        row("9", "updatedAt", "Date", "8 bytes", "Auto (timestamps)", "Record last update time"),
    ],
    "Bookings Collection (bookings)": [
        row("1", "_id", "ObjectId", "24 hex chars", "PK, Auto-generated", "Unique booking identifier"),
        row("2", "user", "ObjectId", "24 hex chars", "FK -> users._id", "User reference"),
        row("3", "stall", "ObjectId", "24 hex chars", "FK -> stalls._id", "Stall reference"),
        row("4", "dome", "ObjectId", "24 hex chars", "FK -> domes._id", "Dome reference"),
        row("5", "stallPrice", "Number", "8 bytes", "Default 0", "Base stall price"),
        row("6", "amount", "Number", "8 bytes", "Optional", "Booking amount"),
        row("7", "defaultMaterials", "Array<Object>", "Variable", "Default []", "Included material snapshot"),
        row("8", "extraMaterials", "Array<Object>", "Variable", "Default []", "Extra material snapshot"),
        row("9", "extraMaterialTotal", "Number", "8 bytes", "Default 0", "Total extra material cost"),
        row("10", "extraMaterialShare", "Number", "8 bytes", "Default 0", "Per booking distributed share"),
        row("11", "grandTotal", "Number", "8 bytes", "Default 0", "Final total amount"),
        row("12", "aadhaarVerification", "ObjectId", "24 hex chars", "Required, FK -> aadhaarverifications._id", "Aadhaar verification reference"),
        row("13", "status", "String", "Up to 9 chars", "Enum: PENDING,APPROVED,PAID,REJECTED,CANCELLED; Default PENDING", "Booking lifecycle status"),
        row("14", "paymentOrderId", "String", "Variable", "Trim, Default empty", "Razorpay order ID"),
        row("15", "paymentId", "String", "Variable", "Trim, Default empty", "Razorpay payment ID"),
        row("16", "paymentSignature", "String", "Variable", "Trim, Default empty", "Payment signature"),
        row("17", "paidAt", "Date", "8 bytes", "Default null", "Payment timestamp"),
        row("18", "createdAt", "Date", "8 bytes", "Auto (timestamps)", "Record creation time"),
        row("19", "updatedAt", "Date", "8 bytes", "Auto (timestamps)", "Record last update time"),
    ],
    "Embedded Booking Material Structure": [
        row("1", "materialId", "ObjectId", "24 hex chars", "FK -> materials._id, Default null", "Material reference"),
        row("2", "name", "String", "Variable", "Required, Trim", "Material name"),
        row("3", "price", "Number", "8 bytes", "Default 0, Min 0", "Unit price"),
        row("4", "quantity", "Number", "8 bytes", "Required, Min 1", "Selected quantity"),
        row("5", "subtotal", "Number", "8 bytes", "Default 0, Min 0", "price x quantity"),
    ],
}


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = "Trade Fair Stall Booking System"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True

    subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.4), Inches(12.5), Inches(1.0))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    p2 = subtitle_tf.paragraphs[0]
    p2.text = f"Data Dictionary\nGenerated on {datetime.now():%d-%m-%Y}"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)


def format_cell_text(cell, value, font_size=11, bold=False, align=PP_ALIGN.LEFT):
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(value)
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold


def add_dictionary_slide(prs, title, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(12.9), Inches(0.5))
    title_tf = title_box.text_frame
    title_tf.clear()
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(20)
    title_p.font.bold = True

    table_shape = slide.shapes.add_table(len(rows) + 1, len(HEADERS), Inches(0.2), Inches(0.7), Inches(12.9), Inches(6.6))
    table = table_shape.table

    table.columns[0].width = Inches(0.8)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(1.2)
    table.columns[3].width = Inches(1.2)
    table.columns[4].width = Inches(2.8)
    table.columns[5].width = Inches(5.3)

    for idx, header in enumerate(HEADERS):
        cell = table.cell(0, idx)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(230, 230, 230)
        format_cell_text(cell, header, font_size=11, bold=True, align=PP_ALIGN.CENTER)

    body_font = 10 if len(rows) <= 10 else 9
    for r_idx, row_data in enumerate(rows, start=1):
        for c_idx, value in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            align = PP_ALIGN.CENTER if c_idx in (0, 2, 3) else PP_ALIGN.LEFT
            format_cell_text(cell, value, font_size=body_font, align=align)


def chunk_rows(items, size):
    for i in range(0, len(items), size):
        yield items[i : i + size]


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    rows_per_slide = 12
    for collection_name, rows in DATA_DICTIONARY.items():
        parts = list(chunk_rows(rows, rows_per_slide))
        for part_index, part_rows in enumerate(parts, start=1):
            if len(parts) > 1:
                title = f"Data Dictionary - {collection_name} (Part {part_index}/{len(parts)})"
            else:
                title = f"Data Dictionary - {collection_name}"
            add_dictionary_slide(prs, title, part_rows)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build_presentation()
    print(f"Created: {output}")
